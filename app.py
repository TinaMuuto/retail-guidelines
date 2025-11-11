import io
import re
import time
from typing import Dict, List, Optional, Tuple
from pathlib import Path

import pandas as pd
import requests
import streamlit as st
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches

# ---------------------- Constants ----------------------
TEMPLATE_PATH = Path("input-template.pptx")
# Important: Ensure these links are publicly accessible (pub?output=csv exports)
DEFAULT_MASTER_URL = "https://docs.google.com/spreadsheets/d/e/2PACX-1vRdNwE1Q_aG3BntCZZPRIOgXEFJ5AHJxHmRgirMx2FJqfttgCZ8on-j1vzxM-muTTvtAHwc-ovDV1qF/pub?output=csv"
DEFAULT_MAPPING_URL = "https://docs.google.com/spreadsheets/d/e/2PACX-1vQPRmVmc0LYISduQdJyfz-X3LJlxiEDCNwW53LhFsWp5fFDS8V669rCd9VGoygBZSAZXeSNZ5fquPen/pub?output=csv"
OUTPUT_NAME = "Muuto_Settings.pptx"

MAX_OVERVIEW_IMAGES = 12
HTTP_TIMEOUT = 10
HTTP_RETRIES = 1
MAX_IMAGE_PX = 1400
JPEG_QUALITY = 85

# ---------------------- Helper Functions ----------------------
def clean_name(name: str) -> str:
    """Cleans shape names for code matching (removes braces and whitespace)."""
    if name is None:
        return ""
    name = name.strip()
    # Removes {{ }} at the start/end
    name = re.sub(r"^\{\{|\}\}$", "", name).strip()
    # Removes whitespace and converts to lowercase for uniformity
    return re.sub(r"\s+", "", name).lower()

def first_run_or_none(shape):
    """Finds the first 'run' in a shape to preserve formatting."""
    try:
        tf = shape.text_frame
        if tf and tf.paragraphs and tf.paragraphs[0].runs:
            return tf.paragraphs[0].runs[0]
    except Exception:
        return None
    return None

def set_text_preserve_format(shape, text: str):
    """Sets text in a shape while trying to preserve the original formatting."""
    try:
        if hasattr(shape, "text_frame") and shape.text_frame:
            run0 = first_run_or_none(shape)
            if run0:
                run0.text = text
            else:
                shape.text_frame.text = text
    except Exception:
        pass

def build_shape_map(slide) -> Dict[str, list]:
    """Creates a dictionary mapping cleaned shape names to shapes."""
    mapping: Dict[str, List] = {}
    for shape in slide.shapes:
        try:
            # Uses clean_name to find matching names from the template
            nm = clean_name(getattr(shape, "name", ""))
            if nm:
                mapping.setdefault(nm, []).append(shape)
        except Exception:
            continue
    return mapping

def safe_find_shape(shape_map: Dict[str, list], key: str, index: int = 0) -> Optional[object]:
    """Helper function to safely find a shape by its cleaned name."""
    clean_key = clean_name(key)
    if clean_key in shape_map and len(shape_map[clean_key]) > index:
        return shape_map[clean_key][index]
    
    # Try alternate names if the primary key fails (for flexibility with user's template names)
    if clean_key.startswith("productpackshot"):
        # If 'productpackshotX' failed, try 'packshotX' (common user simplification)
        alt_key = clean_key.replace("productpackshot", "packshot")
        if alt_key in shape_map and len(shape_map[alt_key]) > index:
            return shape_map[alt_key][index]
            
    # Specific check for Linedrawing typo found in user's template (LlineDrawing)
    if clean_key == "linedrawing":
        alt_key = "llinedrawing"
        if alt_key in shape_map and len(shape_map[alt_key]) > index:
            return shape_map[alt_key][index]
            
    return None

def http_get_bytes(url: str) -> Optional[bytes]:
    """Fetches the content of a URL with retries."""
    if not url:
        return None
    for attempt in range(HTTP_RETRIES + 1):
        try:
            resp = requests.get(url, timeout=HTTP_TIMEOUT, allow_redirects=True)
            if resp.status_code == 200 and resp.content:
                return resp.content
        except Exception:
            pass
        time.sleep(0.2 * attempt)
    return None

def parse_csv_flex(buf: bytes) -> pd.DataFrame:
    """Attempts to parse CSV data with different separators and encodings."""
    if buf is None:
        return pd.DataFrame()
    candidates = [
        {"sep": ";", "encoding": "utf-8-sig"}, # Prioritizes semicolon and UTF-8
        {"sep": ",", "encoding": "utf-8"},      
        {"sep": ";", "encoding": "utf-8"},
        {"sep": "\t", "encoding": "utf-8"},
        {"sep": ",", "encoding": "utf-8-sig"},
        {"sep": ";", "encoding": "latin-1"},
    ]
    for c in candidates:
        try:
            return pd.read_csv(io.BytesIO(buf), sep=c["sep"], encoding=c["encoding"], skipinitialspace=True)
        except Exception:
            continue
    return pd.DataFrame()

def group_key_from_filename(name: str) -> Tuple[str, str]:
    """Extracts group key and file type from the filename."""
    base = Path(name).stem
    lname = base.lower()
    if "floorplan" in lname:
        t = "floorplan"
    elif "linedrawing" in lname or "line_drawing" in lname or "line drawing" in lname:
        t = "linedrawing"
    else:
        ext = Path(name).suffix.lower()
        if ext == ".csv":
            t = "csv"
        elif ext in [".jpg", ".jpeg", ".png"]:
            t = "render"
        else:
            t = "other"
    # Find the key by taking everything AFTER the first " - "
    if " - " in base:
        key = base.split(" - ", 1)[1]
    # Fallback: find the key by taking the last part split by [- _]
    else:
        parts = re.split(r"[-_]", base)
        key = parts[-1] if parts else base
    # Clean up the key by removing type markers (case-insensitive)
    key = re.sub(r"\s+(floorplan|line\s*drawing|linedrawing)$", "", key, flags=re.IGNORECASE).strip()
    return key, t

def base_before_dash(s: str) -> str:
    """Finds the base article number (everything before the first dash)."""
    if not isinstance(s, str):
        s = str(s) if pd.notna(s) else ""
    return s.split("-")[0].strip()

def find_layout_by_name(prs: Presentation, target: str):
    """Finds a layout by name (clean match)."""
    t = clean_name(target)
    for layout in prs.slide_layouts:
        if clean_name(layout.name) == t:
            return layout
    for layout in prs.slide_layouts:
        if t in clean_name(layout.name):
            return layout
    return None

def ensure_presentation_from_path(path: Path) -> Presentation:
    """Ensures the template file exists and can be loaded."""
    if not path.exists():
        raise FileNotFoundError(f"Template not found: {path}")
    return Presentation(str(path))

def load_remote_csv(url: str) -> pd.DataFrame:
    """Fetches and normalizes CSV from a remote URL."""
    content = http_get_bytes(url)
    if content is None:
        return pd.DataFrame()
    df = parse_csv_flex(content)
    return df

def normalize_master(df: pd.DataFrame) -> pd.DataFrame:
    """Normalizes Master Data CSV."""
    if df is None or df.empty:
        return pd.DataFrame(columns=["ITEM NO.", "IMAGE"])
    
    # Aggressive stripping of column names
    cols = {c: c.strip() for c in df.columns}
    df = df.rename(columns=cols)
    
    img_col = next((c for c in df.columns if c.upper() in ["IMAGE URL", "IMAGE DOWNLOAD LINK"] or ("image" in c.lower() and ("url" in c.lower() or "download" in c.lower()))), None)
    item_col = next((c for c in df.columns if c.strip().upper() == "ITEM NO." or ("item" in c.lower() and "no" in c.lower())), None)

    if item_col is None or img_col is None:
        return pd.DataFrame(columns=["ITEM NO.", "IMAGE"])
    
    out = df[[item_col, img_col]].copy()
    out.columns = ["ITEM NO.", "IMAGE"]
    out["ITEM NO."] = out["ITEM NO."].astype(str).str.strip() # Ensure data is stripped
    out["IMAGE"] = out["IMAGE"].astype(str).str.strip()       # Ensure data is stripped
    out["ITEM BASE"] = out["ITEM NO."].apply(base_before_dash)
    return out

def normalize_mapping(df: pd.DataFrame) -> pd.DataFrame:
    """Normalizes Mapping Data CSV."""
    if df is None or df.empty:
        return pd.DataFrame(columns=["OLD Item-variant", "Description", "New Item No."])
    
    # Aggressive stripping of column names
    cols = {c: c.strip() for c in df.columns}
    df = df.rename(columns=cols)
    
    col_old = next((c for c in df.columns if c.lower().strip() in ["old item-variant", "old item variant", "olditem-variant"] or ("old" in c.lower() and "variant" in c.lower())), None)
    col_new = next((c for c in df.columns if c.lower().strip() in ["new item no.", "new item no", "new item number"] or ("new" in c.lower() and ("no" in c.lower() or "number" in c.lower()))), None)
    col_desc = next((c for c in df.columns if c.lower().strip() == "description" or "desc" in c.lower()), None)
    
    if not col_old or not col_new:
        return pd.DataFrame(columns=["OLD Item-variant", "Description", "New Item No."])
    
    if col_desc is None:
        df["__desc__"] = ""
        col_desc = "__desc__"
        
    out = df[[col_old, col_desc, col_new]].copy()
    out.columns = ["OLD Item-variant", "Description", "New Item No."]
    
    out["OLD Item-variant"] = out["OLD Item-variant"].astype(str).str.strip() # Ensure data is stripped
    out["New Item No."] = out["New Item No."].astype(str).str.strip() # Ensure data is stripped
    
    out["OLD BASE"] = out["OLD Item-variant"].apply(base_before_dash)
    out["NEW BASE"] = out["New Item No."].apply(base_before_dash)
    return out

def normalize_pcon(df: pd.DataFrame) -> pd.DataFrame:
    """Normalizes pCon CSV (input file)."""
    if df is None or df.empty:
        return pd.DataFrame(columns=["ARTICLE_NO", "Quantity"])
    
    norm = {c: re.sub(r"[^a-z0-9]", "", c.lower()) for c in df.columns}
    
    article_col = next((c for c in df.columns if norm[c] in {"articleno","article","articlenumber","artno","artnr","artnumber","itemno","itemnumber","articlecode"} or ("article" in norm[c] and "no" in norm[c]) or ("item" in norm[c] and "no" in norm[c])), None)

    qty_col = next((c for c in df.columns if norm[c] in {"qty","quantity","quantities","qtytotal","qtysum"} or "qty" in norm[c]), None)

    if article_col is None:
        return pd.DataFrame(columns=["ARTICLE_NO", "Quantity"])
    
    out = pd.DataFrame()
    out["ARTICLE_NO"] = df[article_col].astype(str).fillna("").str.strip()
    
    if qty_col is not None:
        out["Quantity"] = pd.to_numeric(df[qty_col], errors="coerce").fillna(1).astype(int)
    else:
        out["Quantity"] = 1
        
    out["ARTICLE_BASE"] = out["ARTICLE_NO"].apply(base_before_dash)
    return out[["ARTICLE_NO", "Quantity", "ARTICLE_BASE"]]

def find_packshot_url(article_no: str, mapping_df: pd.DataFrame, master_df: pd.DataFrame) -> Optional[str]:
    """Finds Packshot URL based on article number (incl. mapping)."""
    if master_df is None or master_df.empty:
        return None
    
    if mapping_df is not None and not mapping_df.empty:
        row = mapping_df[mapping_df["OLD Item-variant"] == str(article_no)]
        if row.empty:
            row = mapping_df[mapping_df["OLD BASE"] == base_before_dash(article_no)]
        
        if not row.empty:
            new_item = row.iloc[0]["New Item No."]
            if pd.notna(new_item):
                m = master_df[master_df["ITEM NO."] == str(new_item)]
                if m.empty:
                    m = master_df[master_df["ITEM BASE"] == base_before_dash(str(new_item))]
                if not m.empty:
                    return m.iloc[0]["IMAGE"]

    m = master_df[master_df["ITEM NO."] == str(article_no)]
    if m.empty:
        m = master_df[master_df["ITEM BASE"] == base_before_dash(str(article_no))]
        
    if not m.empty:
        return m.iloc[0]["IMAGE"]
        
    return None

def find_description(article_no: str, mapping_df: pd.DataFrame) -> str:
    """Finds product description from mapping data."""
    if mapping_df is None or mapping_df.empty:
        return ""
        
    row = mapping_df[mapping_df["OLD Item-variant"] == str(article_no)]
    if row.empty:
        row = mapping_df[mapping_df["OLD BASE"] == base_before_dash(article_no)]
        
    if not row.empty:
        desc = row.iloc[0]["Description"]
        return "" if pd.isna(desc) else str(desc).strip()
    return ""

def find_new_item(article_no: str, mapping_df: pd.DataFrame) -> Optional[str]:
    """Finds the new article number from mapping data."""
    if mapping_df is None or mapping_df.empty:
        return None
        
    row = mapping_df[mapping_df["OLD Item-variant"] == str(article_no)]
    if row.empty:
        row = mapping_df[mapping_df["OLD BASE"] == base_before_dash(article_no)]
        
    if not row.empty:
        val = row.iloc[0]["New Item No."]
        return None if pd.isna(val) or val == "" or val.lower() == "nan" else str(val).strip()
    return None

def chunk(lst, n):
    """Splits a list into chunks of size n."""
    for i in range(0, len(lst), n):
        yield lst[i:i+n]

def get_blank_layout(prs: Presentation):
    """Retrieves a blank or empty layout."""
    for layout in prs.slide_layouts:
        if clean_name(layout.name) in ("blank", "empty"):
            return layout
    return prs.slide_layouts[0]

def inherit_layout_names(slide):
    """Copies placeholder names from layout to newly created slide shapes."""
    try:
        layout = slide.slide_layout
        lnames = {}
        for sh in getattr(layout, "shapes", []):
            phf = getattr(sh, "placeholder_format", None)
            if phf is not None:
                idx = getattr(phf, "idx", None)
                nm = getattr(sh, "name", None)
                if idx is not None and nm:
                    lnames[idx] = nm
        for sh in slide.shapes:
            phf = getattr(sh, "placeholder_format", None)
            if phf is not None and getattr(phf, "idx", None) in lnames:
                if not getattr(sh, "name", "") or "placeholder" in sh.name.lower():
                    sh.name = lnames[phf.idx]
    except Exception:
        pass

# --------- Image Helpers ----------
def add_picture_contain(slide, shape, image_bytes: bytes):
    """Inserts an image into a shape, ensuring it fits (contain-fit) without cropping."""
    try:
        if not image_bytes:
            return
        with Image.open(io.BytesIO(image_bytes)) as im:
            im = im.convert("RGB")
            w, h = im.size
            max_dim = min(MAX_IMAGE_PX, max(w, h))
            scale_src_cap = min(1.0, max_dim / float(max(w, h)))
            if scale_src_cap < 1.0:
                im = im.resize((int(w * scale_src_cap), int(h * scale_src_cap)), Image.Resampling.LANCZOS)

            frame_w = int(shape.width)
            frame_h = int(shape.height)
            s = min(frame_w / im.width, frame_h / im.height)
            s = min(s, 1.0)
            target_w = max(1, int(im.width * s))
            target_h = max(1, int(im.height * s))

            buf = io.BytesIO()
            im.resize((target_w, target_h), Image.Resampling.LANCZOS).save(buf, format="JPEG", quality=JPEG_QUALITY, optimize=True)
            buf.seek(0)

            left = shape.left + int((shape.width - target_w) / 2)
            top = shape.top + int((shape.height - target_h) / 2)
            slide.shapes.add_picture(buf, left, top, width=target_w, height=target_h)
            
            # Delete the original anchor shape (if it's a regular shape)
            try:
                if not getattr(shape, "is_placeholder", False):
                    shape.element.getparent().remove(shape.element)
            except Exception:
                pass
            
    except Exception:
        return

def add_picture_into_shape(slide, shape, image_bytes: bytes):
    """
    Forces robust contain-fit to avoid issues with template placeholders.
    (Ignores the advanced insert_picture logic for maximum robustness)
    """
    if not image_bytes or shape is None:
        return
    
    # Go directly to robust contain-fit
    add_picture_contain(slide, shape, image_bytes)


def add_table(slide, anchor_shape, rows: int, cols: int):
    """Creates a table using the anchor shape's position and size."""
    try:
        left = getattr(anchor_shape, 'left', Inches(0.5))
        top = getattr(anchor_shape, 'top', Inches(1.2))
        width = getattr(anchor_shape, 'width', Inches(9.0))
        height = getattr(anchor_shape, 'height', Inches(5.0))
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        return table
    except Exception:
        return None

# --------- Fallback Slide Builders (Used only if Layouts are missing named shapes) ----------

def create_overview_slide_fallback(prs: Presentation, images_batch):
    # Fallback to create a slide if the 'Overview' Layout is inadequate
    slide = prs.slides.add_slide(get_blank_layout(prs))
    inherit_layout_names(slide)
    cols, rows = 4, 3
    margin_x, margin_y = Inches(0.5), Inches(1.0)
    cell_w = Inches(2.0)
    cell_h = Inches(1.5)
    for idx, img_bytes in enumerate(images_batch, start=1):
        r = (idx-1) // cols
        c = (idx-1) % cols
        left = margin_x + c * (cell_w + Inches(0.2))
        top = margin_y + r * (cell_h + Inches(0.2))
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, cell_w, cell_h)
        rect.name = f"Rendering{idx}"
        add_picture_contain(slide, rect, img_bytes)

def create_setting_slide_fallback(prs: Presentation,
                                  group_name: str,
                                  render_bytes: Optional[bytes],
                                  floorplan_bytes: Optional[bytes],
                                  products_df: pd.DataFrame,
                                  mapping_df: pd.DataFrame,
                                  master_df: pd.DataFrame):
    # Fallback to create a slide if the 'Setting' Layout is inadequate
    slide = prs.slides.add_slide(get_blank_layout(prs))
    inherit_layout_names(slide)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.0), Inches(0.6))
    title.name = "SETTINGNAME"
    set_text_preserve_format(title, group_name)

    render_anchor = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5))
    render_anchor.name = "Rendering"
    if render_bytes:
        add_picture_contain(slide, render_anchor, render_bytes)

    line_anchor = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.2), Inches(1.2), Inches(3.0), Inches(3.5))
    line_anchor.name = "Linedrawing"
    if floorplan_bytes:
        add_picture_contain(slide, line_anchor, floorplan_bytes)

    start_top = Inches(5.0)
    cell_w = Inches(1.6)
    cell_h = Inches(1.2)
    gap = Inches(0.2)
    subset = products_df.head(12).copy() if len(products_df) > 12 else products_df.copy()
    for i, row in enumerate(subset.itertuples(index=False), start=1):
        r = (i-1) // 6
        c = (i-1) % 6
        left = Inches(0.5) + c * (cell_w + gap)
        top = start_top + r * (cell_h + Inches(0.6))
        
        pack_anchor = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, cell_w, cell_h)
        pack_anchor.name = f"ProductPackshot{i}"
        pack_url = find_packshot_url(row.ARTICLE_NO, mapping_df, master_df)
        img_bytes = http_get_bytes(pack_url) if pack_url else None
        if img_bytes:
            add_picture_contain(slide, pack_anchor, img_bytes)
        else:
             desc_text = pack_anchor.text_frame.paragraphs[0]
             desc_text.text = f"Image missing for {row.ARTICLE_NO}"
        
        desc_box = slide.shapes.add_textbox(left, top + cell_h + Inches(0.05), cell_w, Inches(0.4))
        desc_box.name = f"PRODUCT DESCRIPTION {i}"
        set_text_preserve_format(desc_box, find_description(row.ARTICLE_NO, mapping_df))

def create_productlist_slide_fallback(prs: Presentation,
                                      group_name: str,
                                      products_df: pd.DataFrame,
                                      mapping_df: pd.DataFrame):
    # Fallback to create a slide if the 'ProductListBlank' Layout is inadequate
    slide = prs.slides.add_slide(get_blank_layout(prs))
    inherit_layout_names(slide)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6))
    set_text_preserve_format(title, f"Products – {group_name}")
    
    # Dummy anchor position used in fallback
    rows = max(1, len(products_df)) + 1
    cols = 3
    table = add_table(slide, None, rows, cols) 
    
    if table is None:
        return
        
    table.cell(0, 0).text = "Quantity"
    table.cell(0, 1).text = "Description"
    table.cell(0, 2).text = "Article No. / New Item No."
    
    r = 1
    for row in products_df.itertuples(index=False):
        table.cell(r, 0).text = str(int(row.Quantity))
        desc = find_description(row.ARTICLE_NO, mapping_df)
        table.cell(r, 1).text = desc
        new_item = find_new_item(row.ARTICLE_NO, mapping_df)
        
        article_text = f"{row.ARTICLE_NO} / {new_item}" if new_item else f"{row.ARTICLE_NO}"
        table.cell(r, 2).text = article_text
        r += 1


# --------- Main Slide Builders (Uses Template Placeholders) ----------
def build_overview_slides(prs: Presentation, overview_layout, rendering_bytes_list: List[bytes]):
    """Builds overview slides using named placeholders in the template."""
    for batch in chunk(rendering_bytes_list, MAX_OVERVIEW_IMAGES):
        slide = prs.slides.add_slide(overview_layout)
        inherit_layout_names(slide)
        shape_map = build_shape_map(slide)
        
        for idx, img_bytes in enumerate(batch, start=1):
            if not img_bytes: continue
            
            pic_key = clean_name(f"Rendering{idx}")
            target_shape = safe_find_shape(shape_map, pic_key)

            if target_shape:
                 add_picture_into_shape(slide, target_shape, img_bytes)

def build_setting_slide(prs: Presentation,
                        setting_layout,
                        group_name: str,
                        render_bytes: Optional[bytes],
                        floorplan_bytes: Optional[bytes],
                        products_df: pd.DataFrame,
                        mapping_df: pd.DataFrame,
                        master_df: pd.DataFrame):
    """Builds individual setting slides using named placeholders in the template."""
    slide = prs.slides.add_slide(setting_layout)
    inherit_layout_names(slide)
    shape_map = build_shape_map(slide)
    
    # 1. Title
    title_shape = safe_find_shape(shape_map, "SETTINGNAME")
    if title_shape:
        set_text_preserve_format(title_shape, group_name)
        
    # 2. Images
    render_shape = safe_find_shape(shape_map, "Rendering")
    if render_shape and render_bytes:
        add_picture_into_shape(slide, render_shape, render_bytes)
        
    # Linedrawing check: safe_find_shape now handles the 'LlineDrawing' typo from your template
    floorplan_shape = safe_find_shape(shape_map, "Linedrawing") 
    if floorplan_shape and floorplan_bytes:
        add_picture_into_shape(slide, floorplan_shape, floorplan_bytes)
        
    # 3. Products and Descriptions
    subset = products_df.head(12).copy() if len(products_df) > 12 else products_df.copy()
    for i, row in enumerate(subset.itertuples(index=False), start=1):
        article_no = row.ARTICLE_NO
        pack_url = find_packshot_url(article_no, mapping_df, master_df)
        img_bytes = http_get_bytes(pack_url) if pack_url else None
        
        # a) Packshot image (safe_find_shape handles the PackshotX vs ProductPackshotX flexibility)
        pic_key = clean_name(f"ProductPackshot{i}")
        pic_shape = safe_find_shape(shape_map, pic_key)
        
        if pic_shape and img_bytes:
            add_picture_into_shape(slide, pic_shape, img_bytes)
        elif pic_shape:
            # Insert a note if the image is missing
            set_text_preserve_format(pic_shape, f"Image for {article_no} missing")
            
        # b) Product Description
        desc_key = clean_name(f"PRODUCT DESCRIPTION {i}")
        desc_shape = safe_find_shape(shape_map, desc_key)
        
        if desc_shape:
            desc = find_description(article_no, mapping_df)
            set_text_preserve_format(desc_shape, desc)

def build_productlist_slide(prs: Presentation,
                            layout,
                            group_name: str,
                            products_df: pd.DataFrame,
                            mapping_df: pd.DataFrame):
    """Builds the product list slide with a table."""
    slide = prs.slides.add_slide(layout)
    inherit_layout_names(slide)
    shape_map = build_shape_map(slide)
    
    # 1. Title
    title_shape = safe_find_shape(shape_map, "Title")
    if title_shape:
        set_text_preserve_format(title_shape, f"Products – {group_name}")
        
    # 2. Table Anchor
    anchor = safe_find_shape(shape_map, "TableAnchor")
    
    if not anchor:
        st.warning(f"WARNING: 'TableAnchor' missing in '{layout.name}'. Using default position.")
    
    rows = max(1, len(products_df)) + 1
    cols = 3
    table = add_table(slide, anchor, rows, cols)
    
    if table is None:
        return
        
    # 3. Fill Table
    table.cell(0, 0).text = "Quantity"
    table.cell(0, 1).text = "Description"
    table.cell(0, 2).text = "Article No. / New Item No."
    
    r = 1
    for row in products_df.itertuples(index=False):
        table.cell(r, 0).text = str(int(row.Quantity))
        desc = find_description(row.ARTICLE_NO, mapping_df)
        table.cell(r, 1).text = desc
        new_item = find_new_item(row.ARTICLE_NO, mapping_df)
        
        article_text = f"{row.ARTICLE_NO} / {new_item}" if new_item else f"{row.ARTICLE_NO}"
        table.cell(r, 2).text = article_text
        r += 1


def layout_has_expected(layout, keys: List[str]) -> bool:
    """Checks if the layout contains the expected placeholders."""
    try:
        names = [clean_name(getattr(sh, "name", "")) for sh in layout.shapes]
    except Exception:
        names = []
        
    # Check for exact match
    if any(clean_name(k) in names for k in keys):
        return True
        
    # Specific check for Setting Layout: Checks if it contains Rendering OR Packshot1
    if clean_name(layout.name) == clean_name("Setting"):
        if any(n.startswith("productpackshot") for n in names) or any(n.startswith("packshot") for n in names) or clean_name("Rendering") in names:
            return True
            
    return False

def preflight_checks() -> Dict[str, str]:
    """Performs preflight checks for template and remote CSVs."""
    results = {}
    try:
        if not TEMPLATE_PATH.exists():
            results["template"] = "Template not found (input-template.pptx)."
        else:
            _ = Presentation(str(TEMPLATE_PATH))
            results["template"] = "OK"
    except Exception as e:
        results["template"] = f"Template unreadable or not a valid .pptx: {e}"
        
    try:
        m = http_get_bytes(DEFAULT_MASTER_URL)
        results["master_csv"] = "OK" if m else "Unavailable (Could not fetch content)"
    except Exception:
        results["master_csv"] = "Unavailable (HTTP error)"
    try:
        mp = http_get_bytes(DEFAULT_MAPPING_URL)
        results["mapping_csv"] = "OK" if mp else "Unavailable (Could not fetch content)"
    except Exception:
        results["mapping_csv"] = "Unavailable (HTTP error)"
    return results

def build_groups(upload_list: List[Dict]) -> Dict[str, Dict]:
    """Groups uploaded files based on the filename."""
    groups: Dict[str, Dict] = {}
    for item in upload_list:
        name = item["name"]
        b = item["bytes"]
        key, t = group_key_from_filename(name)
        if key not in groups:
            groups[key] = {"name": key, "csv": None, "render": None, "floorplan": None}
        if t == "csv":
            groups[key]["csv"] = b
        elif t == "render":
            if groups[key]["render"] is None:
                groups[key]["render"] = b
        elif t in ["floorplan", "linedrawing"]:
            if groups[key]["floorplan"] is None:
                groups[key]["floorplan"] = b
    return groups

def collect_all_renderings(groups: Dict[str, Dict]) -> List[bytes]:
    """Collects all rendering images from the groups."""
    return [g["render"] for g in groups.values() if g.get("render")]

def safe_present(prs: Presentation) -> bytes:
    """Saves the presentation to a bytes buffer."""
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()

# ---------------------- UI ----------------------
st.set_page_config(page_title="Muuto PowerPoint Generator", layout="centered")
st.title("Muuto PowerPoint Generator")
st.write("Upload your group files (CSV and images). The app uses the fixed PowerPoint template and fetches Master Data and Mapping from fixed URLs.")

if "uploads" not in st.session_state:
    st.session_state.uploads = []
if "last_master_df" not in st.session_state:
    st.session_state.last_master_df = None
if "last_mapping_df" not in st.session_state:
    st.session_state.last_mapping_df = None


files = st.file_uploader(
    "User group files (.csv, .jpg, .png). You can add multiple files.",
    type=["csv", "jpg", "jpeg", "png"],
    accept_multiple_files=True,
)

if files:
    existing = {u["name"] for u in st.session_state.uploads}
    for f in files:
        # Check if file is already loaded to avoid reading heavy file multiple times
        if f.name not in existing:
            st.session_state.uploads.append({"name": f.name, "bytes": f.read()})
            existing.add(f.name)

# Single flat file list with remove buttons
if st.session_state.uploads:
    st.subheader("Uploaded Files")
    remove_indices = []
    for idx, item in enumerate(st.session_state.uploads):
        col1, col2 = st.columns([0.9, 0.1])
        with col1:
            size_kb = len(item["bytes"]) / 1024.0
            st.write(f"{item['name']} — {size_kb:.1f}KB")
        with col2:
            if st.button("❌", key=f"rm_{idx}"):
                remove_indices.append(idx)
    if remove_indices:
        for i in sorted(remove_indices, reverse=True):
            del st.session_state.uploads[i]
        st.rerun()

generate = st.button("Generate Presentation")

# ---------------------- Orchestration ----------------------
if generate:
    with st.spinner("Working..."):
        diag = preflight_checks()
        if diag.get("template") != "OK":
            st.error("Template issue: " + diag.get("template", "Unknown"))
        elif not TEMPLATE_PATH.exists():
            st.error("Template file is missing in the repository: input-template.pptx")
        else:
            try:
                # Build groups regardless of whether st.session_state.uploads is empty
                groups = build_groups(st.session_state.uploads)

                if not groups:
                    st.error("Could not form any groups. Please ensure files are uploaded and filenames contain CSV and at least one image/floorplan.")
                    
                # Continue only if groups were successfully formed
                if groups:
                    prs = ensure_presentation_from_path(TEMPLATE_PATH)

                    overview_layout = find_layout_by_name(prs, "Overview") or find_layout_by_name(prs, "Renderings")
                    setting_layout = find_layout_by_name(prs, "Setting")
                    productlist_layout = find_layout_by_name(prs, "ProductListBlank")

                    master_df = normalize_master(load_remote_csv(DEFAULT_MASTER_URL))
                    mapping_df = normalize_mapping(load_remote_csv(DEFAULT_MAPPING_URL))
                    
                    # Save data status for UI feedback
                    st.session_state.last_master_df = master_df
                    st.session_state.last_mapping_df = mapping_df
                    
                    # Check for empty Master/Mapping data
                    if master_df.empty:
                         st.warning("WARNING: Master Data (Image URLs) could not be loaded. Packshots will be missing.")
                    if mapping_df.empty:
                         st.warning("WARNING: Mapping Data (Descriptions/New Article Numbers) could not be loaded. Descriptions and new article numbers will be missing.")
                    

                    # Overview Slides
                    renders = collect_all_renderings(groups)
                    if renders:
                        if overview_layout and layout_has_expected(overview_layout, ["Rendering1"]):
                            build_overview_slides(prs, overview_layout, renders)
                        else:
                            for batch in chunk(renders, MAX_OVERVIEW_IMAGES):
                                create_overview_slide_fallback(prs, batch)

                    # Per group Slides
                    for key in sorted(groups.keys()):
                        g = groups[key]
                        group_name = g["name"]
                        
                        try:
                            pcon_df = normalize_pcon(parse_csv_flex(g["csv"]) if g["csv"] else pd.DataFrame())
                        except Exception:
                            pcon_df = pd.DataFrame(columns=["ARTICLE_NO", "Quantity"])
                        
                        if pcon_df.empty:
                            pcon_df = pd.DataFrame(columns=["ARTICLE_NO", "Quantity"])
                            st.warning(f"WARNING: Could not load product data from CSV for group '{group_name}'.")

                        # Setting Slide (Renders, Linedrawing, Packshots, Descriptions)
                        if setting_layout and layout_has_expected(setting_layout, ["SETTINGNAME", "Rendering", "ProductPackshot1"]):
                            build_setting_slide(prs, setting_layout, group_name, g.get("render"), g.get("floorplan"), pcon_df, mapping_df, master_df)
                        else:
                            create_setting_slide_fallback(prs, group_name, g.get("render"), g.get("floorplan"), pcon_df, mapping_df, master_df)
                            st.info(f"INFO: Setting slide for '{group_name}' used fallback layout (Check placeholder names: SETTINGNAME, Rendering, ProductPackshot1-12).")


                        # Product List Slide (Table)
                        if productlist_layout and layout_has_expected(productlist_layout, ["TableAnchor"]):
                            build_productlist_slide(prs, productlist_layout, group_name, pcon_df, mapping_df)
                        else:
                            create_productlist_slide_fallback(prs, group_name, pcon_df, mapping_df)
                            st.info(f"INFO: Product List slide for '{group_name}' used fallback layout (Check name 'TableAnchor').")

                    ppt_bytes = safe_present(prs)
                    st.success("Your presentation is ready!")
                    st.download_button(
                        "Download Muuto_Settings.pptx",
                        data=ppt_bytes,
                        file_name=OUTPUT_NAME,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
            except Exception as e:
                st.error(f"Something went wrong while generating the presentation: {e}")
                # st.exception(e) # Can be used for deeper debugging

# UI for data status (runs only after a generation, as the URLs are only checked there)
if st.session_state.last_master_df is not None and st.session_state.last_mapping_df is not None:
    st.subheader("Data Connection Status")
    col_m, col_mp = st.columns(2)
    col_m.metric("Master Data Rows", st.session_state.last_master_df.shape[0])
    col_mp.metric("Mapping Data Rows", st.session_state.last_mapping_df.shape[0])
    if st.session_state.last_master_df.empty or st.session_state.last_mapping_df.empty:
        st.warning("WARNING: Zero rows loaded from Master/Mapping CSVs. Check URL accessibility.")
```eof
